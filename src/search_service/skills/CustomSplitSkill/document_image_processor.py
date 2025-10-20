import os
import io
import base64
import tempfile
import concurrent.futures
from PIL import Image
from docx import Document
from pptx import Presentation
from urllib.parse import unquote
import httpx
from .excel_splitter import get_blob_service_client
import time
import subprocess
import platform
from langchain_openai import AzureChatOpenAI
import fitz  
import re
from tenacity import retry, stop_after_attempt, wait_exponential
from functools import partial,lru_cache
import shutil
from io import BytesIO
from .constants import FileExtenstions

# Singleton HTTP client with custom CA certificate
class HttpClient:
    _client = None

    @classmethod
    def get_client(cls):
        if cls._client is None:
            cacert_path = os.path.abspath("cacert.pem")
            print("Path: ", cacert_path)
            cls._client = httpx.Client(verify=cacert_path)
        return cls._client

# Construct the path to the cacert.pem file
cacert_path = os.path.abspath("cacert.pem")
print("Path: ", cacert_path)
# Create an httpx client with the custom CA certificate
client = HttpClient.get_client()
os_type = platform.system()
SOFFICE=r"C:\Program Files\LibreOffice\program\soffice.exe"
# Define the prompts for the AI assistant
doc_image_summarizer_prompt = """
You are an AI assistant for a semiconductor company. Please provide a clear and detailed interpretation of the data or information depicted in the image.
Pay special attention to any text, labels, or regions highlighted in **RED** — these usually indicate key actions, commands, or interface components. Explicitly mention and explain them in your summary. 
Describe the image and its key elements comprehensively. The image could include parts, snapshots of tools, analysis, tables, BOM, engineering drawings, flowcharts, assembly diagrams, etc. 
For tables, include the headers, summarize the content in detail so that any questions asked on that should be answerable. Explain how each row and column are related to each other. Convert the table to a markdown format so that any question asked on top of it across rows and columns can be answered.
For flowcharts, ensure you are explaining the flow of the process, the decision points, and the outcomes accurately and as detailed as possible.
For engineering drawings, explain the components, the assembly process, the dimensions, and any other relevant details.
For BOM, explain the components, the quantities, the part numbers, and any other relevant details.
For analysis, explain the data points, the trends, the insights, and any other relevant details.
For snapshots of tools, assemblies, explain the tool, the purpose, the features, and any other relevant details. Call out the Notes if any in the image.
Avoid bullet points; instead, deliver a coherent, factual summary that captures the essence of the image for analysis. Ensure the description is precise, informative, and leverages any available context from the document. Respond in Markdown format."""

slide_image_summarizer_prompt = """You are an AI Assistant for a semiconductor company. Given an image and context from a presentation slide/pdf document, please provide a clear and detailed interpretation of the data or information presented. Describe the image and its key elements comprehensively. The image could include diagrams, tables, charts, photographs, parts, snapshots of tools, analysis, BOM, engineering drawings, flowcharts, assembly diagrams, etc.
Pay special attention to any text, labels, or regions highlighted in **RED** — these usually indicate key actions, commands, or interface components. Explicitly mention and explain them in your summary.
For tables, include the headers, summarize the content in detail so that any questions asked on that should be answerable. Explain how each row and column are related to each other. Convert the table to a markdown format so that any question asked on top of it across rows and columns can be answered.

For flowcharts, ensure you are explaining the flow of the process, the decision points, and the outcomes accurately and as detailed as possible.

For engineering drawings, explain the components, the assembly process, the dimensions, and any other relevant details.

For BOM, explain the components, the quantities, the part numbers, and any other relevant details.

For analysis, explain the data points, the trends, the insights, and any other relevant details.

For snapshots of tools, assemblies, explain the tool, the purpose, the features, and any other relevant details. Call out the Notes if any in the image.

Explain the slide/page using the image and any available context in a factual and coherent manner. Ensure the explanation is easy to understand for a non-technical audience, but include significant details that are crucial for understanding.

Avoid bullet points; instead, deliver a coherent, factual summary that captures the essence of the image and the slide for analysis. Ensure the description is precise, informative, and leverages any available context from the slide. Respond in Markdown format."""

# Create an instance of the AzureChatOpenAI class
llm = AzureChatOpenAI(
    openai_api_key = os.environ["OPENAI_API_KEY"],
    azure_endpoint = os.environ["AZURE_OPENAI_ENDPOINT"],
    api_version = os.environ["OPENAI_API_VERSION"],
    azure_deployment = os.environ["AZURE_OPENAI_SUMMARIZATION_MODEL"],
    temperature = 0,
    http_client=client
)

number_of_images_to_process = 50

@lru_cache(maxsize=1)
def install_libreoffice_dependencies():
    """
    Installs LibreOffice dependencies if they are not already installed.

    This function checks the operating system and installs LibreOffice using the appropriate method.
    On Linux, it uses the apt-get package manager. On Windows, it checks for the presence of soffice.exe.
    For unsupported operating systems, it raises a runtime error.
    Instruct the user to install it manually if not installed.

    Raises:
        RuntimeError: If LibreOffice is not found on Windows or if the OS is unsupported.
    """
    system = platform.system()
    
    if system == "Linux":
        # Only on Linux do we run apt-get to install LibreOffice
        result = subprocess.run(["which", "libreoffice"], capture_output=True, text=True)
        if result.returncode == 0:
            print("libreoffice is already installed.")
        else:
            subprocess.run(["apt-get", "update"], check=True)
            subprocess.run(["apt-get", "install", "-y", "libreoffice"], check=True)
            print("libreoffice installed successfully.")
    
    elif system == "Windows":
        # On Windows, check for soffice.exe and raise an error if missing
        # If LibreOffice is not found, instruct the user to install it manually
        path = SOFFICE
        if not shutil.which(path):
            raise RuntimeError(
                f"LibreOffice not found at {path}. "
                "Please install it (using GUI, winget, or choco) before running."
            )
    
    else:
        # Raise an error for unsupported operating systems
        raise RuntimeError(f"Unsupported OS: {system}")

def convert_vector_image_to_png(image_part, file_extension):
    """
    Converts a vector image to PNG format using LibreOffice.

    This function takes a vector image and its file extension, writes it to a temporary file,
    and uses LibreOffice to convert it to a PNG file. The path to the converted PNG file is returned.

    Args:
        image_part: An object containing the vector image data. It should have a 'blob' attribute
                    that contains the binary data of the image.
        file_extension: A string representing the file extension of the vector image (e.g., '.svg', '.eps').

    Returns:
        str: The path to the converted PNG file if the conversion is successful, None otherwise.

    Raises:
        Exception: If there is an error during the conversion process, an exception is caught and
                   an error message is printed.
    """
    try:
        print(f"Converting {file_extension.upper()} to PNG using LibreOffice...")
        # Create a temporary file for the input vector image
        with tempfile.NamedTemporaryFile(suffix=file_extension, delete=False) as temp_input:
            temp_input_path = temp_input.name
            temp_input.write(image_part.blob)
        
        print(f"Temporary {file_extension.upper()} input path: {temp_input_path}")
        if not os.path.exists(temp_input_path):
            print(f"Input file {temp_input_path} was not created.")
            return None
        
        # Create a temporary file for the output PNG
        temp_output_path = temp_input_path.replace(file_extension, ".png")
        print(f"Temporary output path: {temp_output_path}")
        
        # Use LibreOffice to convert vector image to PNG
        result = subprocess.run([
            'libreoffice', '--headless', '--convert-to', 'png', temp_input_path, '--outdir', os.path.dirname(temp_output_path)
        ], stdout=subprocess.PIPE, stderr=subprocess.PIPE)
        
        if result.returncode != 0:
            print(f"LibreOffice conversion failed: {result.stderr.decode('utf-8')}")
            return None
        
        if not os.path.exists(temp_output_path):
            print(f"Output file {temp_output_path} was not created.")
            return None
        
        print("Conversion successful.")
        return temp_output_path
    except Exception as e:
        print(f"Failed to convert {file_extension.upper()} to PNG: {str(e)}")
        return None

def convert_to_png_and_base64(image_part, file_extension):
    """
    Converts a vector image to PNG format and encodes it in base64.

    This function first installs the necessary LibreOffice dependencies if they are not already installed.
    It then converts the provided vector image to PNG format using LibreOffice, reads the PNG data,
    and encodes it in base64. The base64-encoded PNG data is returned as a string.

    Args:
        image_part: An object containing the vector image data. It should have a 'blob' attribute
                    that contains the binary data of the image.
        file_extension: A string representing the file extension of the vector image (e.g., '.svg', '.eps').

    Returns:
        str: The base64-encoded PNG data if the conversion and encoding are successful, None otherwise.

    Raises:
        Exception: If there is an error during the conversion or encoding process, an exception is caught and
                   an error message is printed.
    """
    try:
        install_libreoffice_dependencies()
        
        temp_output_path = convert_vector_image_to_png(image_part, file_extension)
        
        if temp_output_path is None:
            return None

        # Read the PNG data
        with open(temp_output_path, "rb") as temp_output:
            image = Image.open(temp_output)
            buffered = io.BytesIO()
            image.save(buffered, format="PNG")
            png_base64 = base64.b64encode(buffered.getvalue()).decode('utf-8')
        
        # Clean up the temporary files
        os.remove(temp_output_path)
        
        return png_base64
    except Exception as e:
        print(f"Failed to convert {file_extension} to PNG and encode to base64: {str(e)}")
        return None

def get_file_from_adls(blob_uri):
    """
    Retrieves a file from Azure Data Lake Storage (ADLS) given a blob URI.
    This function uses environment variables for authentication and connects to the specified
    Azure storage account to download the blob content.
    Args:
        blob_uri (str): The URI of the blob to be retrieved. This URI may contain encoded characters
                        which will be decoded before accessing the blob.
    Returns:
        bytes: The content of the blob as bytes.
    Raises:
        Exception: If there is an issue with accessing the blob or any other error occurs during the process.
    """
    try:
        tenant_id = os.environ["TENET_ID"]
        client_id = os.environ["CLIENT_ID"]
        client_secret = os.environ["CLIENT_SECRET"]
        storage_account_name = os.environ["STORAGE_ACCOUNT_NAME"]
        container_name = "knowledge-mining"

        blob_service_client = get_blob_service_client(
            tenant_id=tenant_id,
            client_id=client_id,
            client_secret=client_secret,
            storage_account_name=storage_account_name,
        )

        # example parent path
        # parent_path = "https://dls2lamuswdfabdev001.blob.core.windows.net/knowledge-mining/commonSpec/Single%20Product%20Spec/ASE/Software%20Spec/C3/Sabre3D/2018/20181002_k21%208%20inch%20LAM%20Plate_Response.doc"
        # This path contains %20 and the url needs to be cleaned
        unquoted_blob_uri = unquote(blob_uri)
        parent_path = unquoted_blob_uri.split(container_name)[1]
        blob_client = blob_service_client.get_blob_client(
            container=container_name, blob=parent_path
        ).download_blob()
        blob_data = blob_client.content_as_bytes()
        return blob_data
    except Exception as e:
        print(f"Failed to retrieve file from ADLS: {str(e)}")
        return None

def get_file_content_as_bytes(file_path):
    """
    Reads the content of a file and returns it as a bytes object.
    Args:
        file_path (str): The path to the input file.
    Returns:
        bytes: The content of the file as a bytes object.
    """
    with open(file_path, 'rb') as file:
        file_content_bytes = file.read()
    return file_content_bytes

def create_temp_file_from_bytes(file_content_bytes, file_extension):
    """
    Creates a temporary file from the provided bytes content and returns the file path.
    Args:
        file_content_bytes (bytes): The content of the file as bytes.
        file_extension (str): The file extension to be used for the temporary file.
    Returns:
        str: The path to the temporary file.
    """
    with tempfile.NamedTemporaryFile(delete=False, suffix=file_extension) as temp_file:
        temp_file.write(file_content_bytes)
        temp_file_path = temp_file.name
    return temp_file_path

def process_image_relationship(rel, os_type):
    """
    Processes an image relationship and converts it to PNG format if necessary.

    This function checks if the relationship is external or if it does not contain an image.
    If either condition is true, it returns None. Otherwise, it processes the image based on its content type.
    For WMF/EMF images, it converts them to PNG format using platform-specific methods.
    For other image types, it converts them to PNG format and encodes them in base64.

    Args:
        rel: An object representing the image relationship. It should have attributes 'is_external',
             'target_ref', and 'target_part'. The 'target_part' should have attributes 'blob' and 'content_type'.
        os_type: A string representing the operating system type (e.g., "Windows", "Linux").

    Returns:
        tuple: A tuple containing the base64-encoded PNG data and the content type if the conversion is successful,
               (None, None) otherwise.

    Raises:
        Exception: If there is an error during the conversion process, an exception is caught and
                   an error message is printed.
    """
    if rel.is_external or "image" not in rel.target_ref:
        return None, None

    image = rel.target_part.blob
    content_type = rel.target_part.content_type

    if content_type in ['image/x-wmf', 'image/x-emf']:
        print("Converting WMF/EMF image to PNG...")
        if os_type == "Windows":
            return convert_wmf_to_png_windows(rel.target_part), content_type
        else:
            return convert_to_png_and_base64(rel.target_part, "." + content_type.split("/")[-1].split("-")[-1]), content_type
        
    with io.BytesIO(image) as image_stream:
        img = Image.open(image_stream)
        with io.BytesIO() as buffered:
            img.save(buffered, format="PNG")
            base64_image = base64.b64encode(buffered.getvalue()).decode('utf-8')
            return base64_image, content_type

def get_base64_images_from_docx(docx_bytes):
    """
    Extracts images from a DOCX file and returns them as base64 encoded strings.
    Args:
        docx_bytes (bytes): The byte content of the DOCX file.
    Returns:
        List[str]: A list of base64 encoded strings representing the images extracted from the DOCX file.
    Raises:
        Exception: If there is an error during image extraction, it will be caught and logged.
    """
    print("Trying to extract images from word file.")
    # create a temporary named file for processing the docx file
    temp_file_path = create_temp_file_from_bytes(docx_bytes, FileExtenstions.DOCX.value)
    print(f"Temporary file created at: {temp_file_path}")
    # Load the DOCX document
    try:
        # print("Temp file path: ", temp_file_path)
        doc = Document(temp_file_path)
    except Exception as e:
        print(f"Failed to load the document: {str(e)}")
        return []
    print(f"Number of images in the document: {len(doc.inline_shapes)}")
    # List to store base64 encoded images
    base64_images = []
    content_type_list =[]
    for rel in doc.part.rels.values():
        try:
            base64_image, content_type = process_image_relationship(rel, os_type)
            if base64_image:
                base64_images.append(base64_image)
                content_type_list.append(content_type)
        except Exception as e:
            print("Unable to extract image: ", str(e))
            continue
    print("Content type list: ", content_type_list)
    os.remove(temp_file_path)
    print("Images extracted from word file with length: ", len(base64_images))
    return base64_images[:number_of_images_to_process]

def convert_wmf_to_png_windows(image_part):
    """
    Converts a WMF image to PNG format on Windows and encodes it in base64.

    This function takes a WMF image, converts it to PNG format using the PIL library, 
    and then encodes the PNG data in base64. The base64-encoded PNG data is returned as a string.

    Args:
        image_part: An object containing the WMF image data. It should have a 'blob' attribute
                    that contains the binary data of the image.

    Returns:
        str: The base64-encoded PNG data if the conversion and encoding are successful, None otherwise.

    Raises:
        Exception: If there is an error during the conversion or encoding process, an exception is caught and
                   an error message is printed.
    """
    try:
        image_data = image_part.blob
        image = Image.open(io.BytesIO(image_data))
        with io.BytesIO() as output:
            image.save(output, format="PNG")
            return base64.b64encode(output.getvalue()).decode('utf-8')
    except Exception as e:
        print(f"Failed to convert WMF to PNG: {str(e)}")
        return None

def get_base64_images_from_pptx(pptx_bytes):
    """
    Extracts images from a PPTX file and returns them as base64 encoded strings along with their slide numbers.
    Args:
        pptx_bytes (bytes): The byte content of the PPTX file.
    Returns:
        list of dict: A list of dictionaries, each containing:
            - 'slide_number' (int): The slide number where the image was found.
            - 'base64_image' (str): The base64 encoded string of the image.
    Raises:
        Exception: If there is an error processing the PPTX file.
    Example:
        pptx_bytes = open('example.pptx', 'rb').read()
        images = get_base64_images_from_pptx(pptx_bytes)
        for image in images:
            print(f"Slide {image['slide_number']}: {image['base64_image']}")
    """
    temp_file_path = create_temp_file_from_bytes(pptx_bytes, ".pptx")
    presentation = load_presentation(temp_file_path)
    if not presentation:
        return []

    images_with_slide_numbers = extract_images_from_presentation(presentation)
    os.remove(temp_file_path)
    return images_with_slide_numbers[:number_of_images_to_process]

def load_presentation(temp_file_path):
    """
    Loads a presentation from a given file path.

    This function attempts to load a presentation from the specified temporary file path using the
    `Presentation` class. It checks if the presentation contains any slides and prints the number
    of slides if they exist. If the presentation contains no slides or if an error occurs during
    loading, it prints an appropriate message and returns None.

    Args:
        temp_file_path (str): The path to the temporary file containing the presentation.

    Returns:
        Presentation: An instance of the `Presentation` class if the presentation is successfully loaded
                      and contains slides, None otherwise.

    Raises:
        Exception: If there is an error during the loading process, an exception is caught and
                   an error message is printed.
    """
    try:
        presentation = Presentation(temp_file_path)
        if presentation.slides is None:
            print("The presentation contains no slides.")
            return None
        else:
            print(f"Number of slides in the presentation: {len(presentation.slides)}")
            return presentation
    except Exception as e:
        print(f"Failed to load the presentation: {str(e)}")
        return None

def extract_images_from_presentation(presentation):
    """
    Extracts images from all slides in a presentation.

    This function iterates through each slide in the given presentation, extracts images from each slide,
    and associates them with their respective slide numbers. The extracted images along with their slide
    numbers are collected into a list and returned.

    Args:
        presentation (Presentation): An instance of the `Presentation` class containing the slides from
                                     which images are to be extracted.

    Returns:
        list: A list of tuples, where each tuple contains an image and its corresponding slide number.

    Raises:
        Exception: If there is an error during the extraction process, an exception is caught and
                   an error message is printed.
    """
    images_with_slide_numbers = []
    for slide_number, slide in enumerate(presentation.slides, start=1):
        print(f"Processing slide {slide_number}...")
        images_with_slide_numbers.extend(extract_images_from_slide(slide, slide_number))
    return images_with_slide_numbers

def extract_images_from_slide(slide, slide_number):
    """
    Extracts images from all slides in a presentation.

    This function iterates through each slide in the given presentation, extracts images from each slide,
    and associates them with their respective slide numbers. The extracted images along with their slide
    numbers are collected into a list and returned.

    Args:
        presentation (Presentation): An instance of the `Presentation` class containing the slides from
                                     which images are to be extracted.

    Returns:
        list: A list of tuples, where each tuple contains an image and its corresponding slide number.

    Raises:
        Exception: If there is an error during the extraction process, an exception is caught and
                   an error message is printed.
    """
    images = []
    for shape in slide.shapes:
        if shape.shape_type == 13:  # 13 corresponds to picture shape type
            image_part = shape.image
            base64_image = process_image_part(image_part, slide_number)
            if base64_image:
                images.append({
                    'slide_number': slide_number,
                    'base64_image': base64_image
                })
    return images

def process_image_part(image_part, slide_number):
    """
    Processes an image part and converts it to base64-encoded PNG format if necessary.

    This function checks the content type of the provided image part. If the image is in WMF or EMF format,
    it converts the image to PNG format using platform-specific methods and encodes it in base64. For other
    image types, it directly converts the image to base64 format. The base64-encoded image data is returned.

    Args:
        image_part: An object containing the image data. It should have attributes 'blob' and 'content_type'.
        slide_number (int): The number of the slide from which the image is being processed.

    Returns:
        str: The base64-encoded image data if the conversion and encoding are successful, None otherwise.

    Raises:
        Exception: If there is an error during the conversion or encoding process, an exception is caught and
                   an error message is printed.
    """
    if image_part.content_type in ["image/x-wmf", "image/x-emf"]:
        print(f"Converting image on {slide_number} from {image_part.content_type} to PNG...")
        if os_type == "Windows":
            return convert_wmf_to_png_windows(image_part)
        else:
            return convert_to_png_and_base64(image_part, "." + image_part.content_type.split("/")[-1].split("-")[-1])
    else:
        return convert_image_to_base64(image_part)

def convert_image_to_base64(image_part):
    """
    Converts an image to base64-encoded PNG format.

    This function takes an image part, reads its binary data, converts it to PNG format using the PIL library,
    and then encodes the PNG data in base64. The base64-encoded PNG data is returned as a string.

    Args:
        image_part: An object containing the image data. It should have a 'blob' attribute
                    that contains the binary data of the image.

    Returns:
        str: The base64-encoded PNG data if the conversion and encoding are successful.

    Raises:
        Exception: If there is an error during the conversion or encoding process, an exception is caught and
                   an error message is printed.
    """
    image_bytes = image_part.blob
    with io.BytesIO(image_bytes) as image_stream:
        img = Image.open(image_stream)
        with io.BytesIO() as buffered:
            img.save(buffered, format="PNG")
            return base64.b64encode(buffered.getvalue()).decode('utf-8')
def convert_ppt_to_pdf(blob_bytes: bytes) -> bytes:
    """
    Converts a .ppt file to a .pdf format using LibreOffice.

    Args:
        blob_bytes (bytes): The byte content of the .ppt file to be converted.

    Returns:
        bytes: The byte content of the converted .pdf file.

    Raises:
        Exception: If there is an error during the conversion process.
    """
    if not os.path.exists(SOFFICE):
        raise Exception("LibreOffice is not installed or the path is incorrect.")

    tmp_in = None
    pdf_path = None

    try:
        # Create a temporary file with the .ppt extension and write the input bytes to it
        with tempfile.NamedTemporaryFile(suffix=".ppt", delete=False) as tmp_in:
            tmp_in.write(blob_bytes)
            tmp_in.flush()  # Ensure all data is written to the file

        # Get the directory of the temporary file to use as the output directory
        out_dir = os.path.dirname(tmp_in.name)

        # Use LibreOffice in headless mode to convert the .ppt to .pdf
        subprocess.run([
            SOFFICE,
            "--headless",
            "--convert-to", "pdf",
            tmp_in.name,
            "--outdir", out_dir
        ], check=True)

        # Read the generated .pdf file back into memory
        pdf_path = os.path.splitext(tmp_in.name)[0] + ".pdf"
        with open(pdf_path, "rb") as f:
            pdf_bytes = f.read()

        # Return the .pdf content as bytes
        return pdf_bytes

    except subprocess.CalledProcessError as cpe:
        raise Exception(f"LibreOffice conversion failed: {cpe}") from cpe
    except Exception as e:
        raise Exception(f"Failed to convert PPT to PDF: {e}") from e
    finally:
        # Clean up temporary files if they exist
        try:
            if tmp_in and os.path.exists(tmp_in.name):
                os.remove(tmp_in.name)
        except Exception:
            pass
        try:
            if pdf_path and os.path.exists(pdf_path):
                os.remove(pdf_path)
        except Exception:
            pass
def get_base64_images_from_pdf(pdf_bytes):
    """
    Extracts images from a PDF file and returns them as base64 encoded strings along with their page numbers.
    Also saves the images to the specified output folder.
    Args:
        pdf_bytes (bytes): The byte content of the PDF file.
        output_folder (str): The path to the folder where the images will be saved.
    Returns:
        list of dict: A list of dictionaries, each containing:
            - 'page_number' (int): The page number where the image was found.
            - 'base64_image' (str): The base64 encoded string of the image.
    Raises:
        Exception: If there is an error processing the PDF file.
    """
    # Create a temporary named file for processing the pdf file
    with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as temp_file:
        temp_file.write(pdf_bytes)
        temp_file_path = temp_file.name
        print(f"Temporary file created at: {temp_file_path}")

    images_with_page_numbers = []

    # Open the PDF document
    try:
        pdf_document = fitz.open(temp_file_path)
    except Exception as e:
        print(f"Failed to load the PDF: {str(e)}")
        return []

    # Iterate through all pages in the PDF
    for page_number in range(len(pdf_document)):
        print(f"Processing page {page_number + 1}...")
        page = pdf_document.load_page(page_number)
        image_list = page.get_images(full=True)

        # Check if any image on the page is larger than 50x50 pixels
        process_page_as_image = False
        for img_index, img in enumerate(image_list):
            xref = img[0]
            base_image = pdf_document.extract_image(xref)
            image_width = base_image["width"]
            image_height = base_image["height"]

            if image_width >= 50 and image_height >= 50:
                process_page_as_image = True
                break

        if process_page_as_image:
            # Render the entire page as an image with scaling
            matrix = fitz.Matrix(2, 2)
            pix = page.get_pixmap(matrix=matrix)
            image_bytes = pix.tobytes(output="png")

            # Convert image to base64
            base64_image = base64.b64encode(image_bytes).decode('utf-8')
            images_with_page_numbers.append({
                'page_number': page_number + 1,
                'base64_image': base64_image
            })

    pdf_document.close()
    os.remove(temp_file_path)
    return images_with_page_numbers[:number_of_images_to_process]

@retry(reraise=True, stop=stop_after_attempt(5), wait=wait_exponential(multiplier=1, min=5, max=45))
def extract_image_summary(image_b64, doc_context=None, document_path=None):
    """
    Extracts a summary from a base64 encoded image using a language model.

    Args:
        image_b64 (str): The base64 encoded image string.
        doc_context (str, optional): Additional context for the document, if available. Defaults to None.

    Returns:
        str: The summary of the image content if successful, None otherwise.

    Raises:
        Exception: If there is an error during the image summary extraction process.
    """
    try:
        print("Extracting Image Summary...", document_path)   
        base64_image = image_b64
        if doc_context==None:
            messages=[
            {"role": "system", "content": doc_image_summarizer_prompt},
            {"role": "user", "content": [
                {"type": "image_url", "image_url": {
                    "url": f"data:image/png;base64,{base64_image}"}
                }
            ]}
            ]
            ai_message = llm.invoke(messages)
        else:
            messages=[
            {"role": "system", "content": slide_image_summarizer_prompt},
            {"role": "user", "content": [
                {"type": "image_url", "image_url": {
                    "url": f"data:image/png;base64,{base64_image}"}
                },
                {"type": "text", "text": "Additional Context: " + doc_context}
            ]}
            ]
            ai_message = llm.invoke(messages)
        return (ai_message.content)
    except Exception as e:
        print(f"Unable to extract Image Summary for document in {document_path}; Error: {str(e)}")
        return None
    
def process_slide_images(image, document_content, document_path):
    """
    Processes an image associated with a slide and extracts a summary based on the document content.

    Args:
        image (dict): A dictionary containing image data, including 'slide_number' and 'base64_image'.
        document_content (list): A list of dictionaries, each containing 'row' and 'content' keys, representing the document's content.

    Returns:
        dict or None: A dictionary with 'row' and 'content' keys if the image summary is successfully extracted, otherwise None.
    """
    slide_number = image['slide_number']
    print(f"Processing image for slide {slide_number}...")
    # Find the content for the corresponding slide number
    # slide_content = next((item['content'] for item in document_content if item['row'] == slide_number), None)
    slide_content = document_content
    if slide_content:
        start_time = time.time()  # Start time
        image_summary = extract_image_summary(image['base64_image'], doc_context=slide_content, document_path=document_path)
        end_time = time.time()  # End time
        latency = end_time - start_time
        print(f"Latency for slide {slide_number}: {latency:.2f} seconds")
        if image_summary is not None:
            return {
                'row': slide_number,
                'content': "Image Summary: " + image_summary
            }
        else:
            print(f"Image summary for slide {slide_number} is None.")
    return None

def process_pdf_images(image, document_content, document_path):
    """
    Processes an image associated with a slide and extracts a summary based on the document content.

    Args:
        image (dict): A dictionary containing image data, including 'page_number' and 'base64_image'.
        document_content (list): A list of dictionaries, each containing 'row' and 'content' keys, representing the document's content.

    Returns:
        dict or None: A dictionary with 'row' and 'content' keys if the image summary is successfully extracted, otherwise None.
    """
    page_number = image['page_number']
    print(f"Processing image for page {page_number}...")
    # Find the content for the corresponding slide number
    # slide_content = next((item['content'] for item in document_content if item['row'] == slide_number), None)
    page_content = document_content
    start_time = time.time()  # Start time
    image_summary = extract_image_summary(image['base64_image'], doc_context=page_content, document_path=document_path)
    end_time = time.time()  # End time
    latency = end_time - start_time
    print(f"Latency for page {page_number}: {latency:.2f} seconds")
    if image_summary is not None:
        return {
            'row': page_number,
            'content': "Image Summary: " + image_summary
        }
    else:
        print(f"Image summary for slide {page_number} is None.")
    return None

def process_images_in_document(document_path_or_bytes, document_content=None):
    """
    Processes images in a given document and extracts summaries for each image.

    Args:
        document_path_or_bytes (str or bytes): Blob URI, local file path, or raw PDF file bytes.
        document_content (list, optional): A list of dicts with 'row' and 'content'.

    Returns:
        list: A list of image summaries (dicts with 'row' and 'content').
    """

    file_content_bytes = None
    document_path_str = None

    # Determine how input was provided
    if isinstance(document_path_or_bytes, bytes):
        file_content_bytes = document_path_or_bytes
        file_extension = ".pdf"  # assume PDF if process_images is being used
    elif isinstance(document_path_or_bytes, str):
        _, ext = os.path.splitext(document_path_or_bytes)
        file_extension = ext.lower()

        if document_path_or_bytes.lower().startswith("http"):
            # Download from ADLS via blob URI
            file_content_bytes = get_file_from_adls(document_path_or_bytes)
        else:
            # Read local file
            file_content_bytes = get_file_content_as_bytes(document_path_or_bytes)

        document_path_str = document_path_or_bytes
    else:
        print("Unsupported input type for process_images_in_document.")
        return []

    if file_extension not in [".pdf", ".docx", ".pptx"]:
        print(f"Unsupported document type for image processing ({file_extension}). Please use a PDF, DOCX, or PPTX file.")
        return []


    if not file_content_bytes:
        print("No file content retrieved, cannot process images.")
        return []

    # Prepare content dictionary for image summarization context
    if document_content and all(isinstance(item, dict) and 'row' in item and 'content' in item for item in document_content):
        document_content_dict = {item['row']: item['content'] for item in document_content}
    else:
        document_content_dict = {}

    temp_file_path = create_temp_file_from_bytes(file_content_bytes, file_extension)
    try:
        if file_extension == ".pdf":
            print("Processing PDF file...")
            results = process_pdf(file_content_bytes, document_content_dict, temp_file_path)

        elif file_extension == ".docx":
            print("Processing DOCX file...")
            results = process_docx(file_content_bytes, temp_file_path)

        elif file_extension == ".pptx":
            print("Processing PPTX file...")
            results = process_pptx(file_content_bytes, document_content_dict, temp_file_path)

    finally:
        # Clean up the temporary file
        if os.path.exists(temp_file_path):
            os.remove(temp_file_path)


    return results

def process_docx(file_content_bytes, document_path):
    print("Processing DOCX file...")
    base64_images = get_base64_images_from_docx(file_content_bytes)
    base64_images = [{'base64_image': img, 'document_path': document_path} for img in base64_images]
    print(f"Number of images in the document: {len(base64_images)}")
    return process_images(base64_images, extract_image_summary)

def process_pptx(file_content_bytes, document_content_dict, document_path):
    print("Processing PPTX file...")
    images_with_slide_numbers = get_base64_images_from_pptx(file_content_bytes)
    print("Length of image list: ", len(images_with_slide_numbers))
    return process_images(images_with_slide_numbers, process_slide_images, document_content_dict, document_path)

def process_pdf(file_content_bytes, document_content_dict, document_path):
    print("Processing PDF file...")
    if file_content_bytes is None:
        print("Getting file content as bytes.")
        file_content_bytes = get_file_content_as_bytes(document_path)
    images_with_page_numbers = get_base64_images_from_pdf(file_content_bytes)
    print("Length of image list: ", len(images_with_page_numbers))
    return process_images(images_with_page_numbers, process_pdf_images, document_content_dict, document_path)

def process_images(images, processing_function, document_content_dict=None, document_path=None):
    image_summaries = []
    with concurrent.futures.ThreadPoolExecutor(max_workers=2) as executor:
        futures = []
        for i, image in enumerate(images):
            if document_content_dict:
                page_or_slide_num = image.get('slide_number') or image.get('page_number') or 1
                futures.append(
                    executor.submit(
                        processing_function,
                        image,
                        document_content_dict.get(page_or_slide_num, document_content_dict.get(1)),
                        document_path
                    )
                )
            else:
                futures.append(
                    executor.submit(
                        partial(processing_function, image['base64_image'], document_path=image.get('document_path'))
                    )
                )

            
            if (i + 1) % 2 == 0:
                time.sleep(0.5)

        for future in concurrent.futures.as_completed(futures):
            result = future.result()
            if result is not None:
                image_summaries.append(result)

    return image_summaries

def extract_document_title(document_text):
    """
    Extracts the title of a document using an AI assistant.
    This function sends the provided document text to an AI assistant, which is 
    instructed to generate a title based on the context of the document. The 
    function returns the title as a string. If an error occurs during the 
    process, an empty string is returned.
    Args:
        document_text (str): The text content of the document from which the title 
                             needs to be extracted.
    Returns:
        str: The extracted title of the document, or an empty string if an error occurs.
    """
    try:
        messages=[
        {"role": "system", "content": "You are an AI assistant for a semiconductor company. Please provide a title for the document from the context passed. If there is document name in the passed context, then return that as title of the document. Return ONLY the title and nothing else."},
        {"role": "user", "content": [
                        {"type": "text", "text": document_text}
        ]}
        ]

        ai_message = llm.invoke(messages)
        return "Document Title: " + (ai_message.content)
    except Exception as e:
        print(f"Unable to extract Document Title; Error: {str(e)}")
        return ""

def convert_document_to_pdf(blob_uri):
    try:
        install_libreoffice_dependencies()
        # Create a temporary pdf file for the input document
        with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as temp_pdf:
            temp_pdf_path = temp_pdf.name
        subprocess.call(["libreoffice", "--headless", "--convert-to", "pdf", "--outdir", os.path.dirname(temp_pdf_path), blob_uri])
        return temp_pdf_path
    except Exception as e:
        print(f"Failed to convert document to PDF: {str(e)}")
        return None
def convert_doc_docx_to_pdf(blob_bytes: bytes) -> BytesIO:
    """
    Converts a document file (either .doc or .docx) to a PDF format using LibreOffice.

    This function takes the bytes of a document file, determines its format, and uses
    LibreOffice in headless mode to convert it to a PDF. The resulting PDF is returned
    as a BytesIO object.

    Args:
        blob_bytes (bytes): The byte content of the document to be converted.

    Returns:
        BytesIO: A BytesIO object containing the PDF data.

    Raises:
        Exception: If there is an error during the conversion process, an exception is raised
                   with a message indicating the failure.
    """
    # Ensure that LibreOffice dependencies are installed and available
    install_libreoffice_dependencies()
    
    # Determine the file extension based on the first two bytes of the input
    # If the bytes match 'PK', assume it's a .docx file; otherwise, assume .doc
    ext = ".docx" if blob_bytes[:2] == b'PK' else ".doc"
    
    # Create a temporary file with the determined extension and write the input bytes to it
    with tempfile.NamedTemporaryFile(suffix=ext, delete=False) as tmp_in:
        tmp_in.write(blob_bytes)
        tmp_in.flush()  # Ensure all data is written to the file
    
    # Get the directory of the temporary file to use as the output directory
    out_dir = os.path.dirname(tmp_in.name)

    # Use LibreOffice in headless mode to convert the document to PDF
    subprocess.run([
        "libreoffice", "--headless",
        "--convert-to", "pdf",
        tmp_in.name,
        "--outdir", out_dir
    ], check=True)

    # Read the generated PDF file back into memory
    pdf_path = os.path.splitext(tmp_in.name)[0] + ".pdf"
    with open(pdf_path, "rb") as f:
        pdf_bytes = f.read()

    # Remove the temporary input and output files to clean up
    os.remove(tmp_in.name)
    os.remove(pdf_path)

    # Return the PDF content as a BytesIO object for in-memory file handling
    return BytesIO(pdf_bytes)

def extract_sections_from_word_document(blob_uri):

    # Retrieve the file content from Azure Data Lake Storage
    file_content_bytes = get_file_from_adls(blob_uri)

    # create a temporary named file for processing the docx file
    with tempfile.NamedTemporaryFile(delete=False, suffix=".docx")\
          as temp_file:
        temp_file.write(file_content_bytes)
        temp_file_path = temp_file.name
        print(f"Temporary file created at: {temp_file_path}")
    
    # Load the DOCX document
    try:
        print("Temp file path: ", temp_file_path)
        doc = Document(temp_file_path)
    except Exception as e:
        print(f"Failed to load the document: {str(e)}")
        return []
    
    try:
        markdown_content=[]
        # Process each paragraph
        for paragraph in doc.paragraphs:
            # Skip empty paragraphs
            if not paragraph.text.strip():
                continue

            text = paragraph.text
            style = paragraph.style.name
            # print("style1", style)
            # Convert heading styles
            if 'heading' in style.lower():
                level = ''.join(re.findall(r'\d+', style))
                markdown_content.append(f"{'#' * int(level)} {text}\n")
        
        return "\n".join(markdown_content) # Return the sections identified along with the section text
    except Exception as e:
        print(f"Failed to extract sections from the document: {str(e)}")
        return ""
    finally:
        os.remove(temp_file_path)